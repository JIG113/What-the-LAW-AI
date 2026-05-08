from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


def _paginate_blocks(blocks: list[str], page_size: int = 12) -> list[str]:
    clean = [b.strip() for b in blocks if b and b.strip()]
    if not clean:
        return [""]
    return ["\n\n".join(clean[i : i + page_size]) for i in range(0, len(clean), page_size)]


def _parse_pdf(path: Path) -> list[str]:
    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        pages.append((page.extract_text() or "").strip())
    return pages or [""]


def _parse_docx(path: Path) -> list[str]:
    doc = DocxDocument(str(path))
    blocks = [p.text for p in doc.paragraphs]
    return _paginate_blocks(blocks)


def _parse_xlsx(path: Path) -> list[str]:
    wb = load_workbook(str(path), data_only=True)
    blocks: list[str] = []
    for ws in wb.worksheets:
        blocks.append(f"[시트] {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if vals:
                blocks.append(" | ".join(vals))
    return _paginate_blocks(blocks)


def _parse_pptx(path: Path) -> list[str]:
    prs = Presentation(str(path))
    pages: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        pages.append(f"[슬라이드 {idx}]\n" + "\n".join(texts))
    return pages or [""]


def parse_file(path: str) -> list[str]:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(p)
    if suffix == ".docx":
        return _parse_docx(p)
    if suffix == ".xlsx":
        return _parse_xlsx(p)
    if suffix == ".pptx":
        return _parse_pptx(p)
    raw = p.read_bytes()
    text = raw.decode("utf-8", errors="ignore")
    return _paginate_blocks(text.split("\n\n"))
